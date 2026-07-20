#!/usr/bin/env python3
"""Standalone one-slide PPTX: n8n learner login details."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

NAVY=RGBColor(0x0B,0x12,0x20); BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72)
LIGHT=RGBColor(0xF5,0xF8,0xFC); WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0)
LBLUE=RGBColor(0xB3,0xD4,0xFF)
FONT="Arial"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def rect(s,x,y,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background(); sh.shadow.inherit=False; return sh

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    if runs and not isinstance(runs[0],list): runs=[runs]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(4); p.space_before=Pt(0)
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz)
            r.font.bold=bold; r.font.color.rgb=col; r.font.name=FONT
    return tb

s=prs.slides.add_slide(BLANK)

# White background
rect(s,0,0,SW,SH,WHITE)

# Top blue bar + left accent bar (matching course style)
rect(s,0,0,SW,Inches(0.22),BLUE)
rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
rect(s,Inches(0.55),Inches(0.60),Inches(0.14),Inches(0.66),BLUE)

# Kicker + Title
txt(s,Inches(0.8),Inches(0.5),Inches(11.8),Inches(0.34),[
    [("COURSE SETUP",14,BLUE,True)]])
txt(s,Inches(0.8),Inches(0.82),Inches(12.0),Inches(0.85),[
    [("Your n8n Login Details",29,INK,True)]])
rect(s,Inches(0.8),Inches(1.55),Inches(12.0),Pt(2),LINE)

# ── URL instruction box ─────────────────────────────────────────
rect(s,Inches(0.85),Inches(1.85),Inches(11.7),Inches(1.05),LIGHT)
txt(s,Inches(1.1),Inches(1.97),Inches(11.2),Inches(0.44),[
    [("URL  ",15,GREY,True),
     ("https://",15,INK,False),
     ("[YOUR EMAIL PREFIX]",15,BLUE,True),
     (".app.n8n.cloud/signin",15,INK,False)]])
txt(s,Inches(1.1),Inches(2.44),Inches(11.2),Inches(0.36),[
    [("e.g.  ",14,GREY,False),
     ("https://n8n1002.app.n8n.cloud/signin",14,TEAL,True)]])

# ── Email list label ────────────────────────────────────────────
txt(s,Inches(0.85),Inches(3.05),Inches(11.7),Inches(0.38),[
    [("Pick ONE email address — one per person:",14,BLUE,True)]])

# ── Email list box — two columns ────────────────────────────────
rect(s,Inches(0.85),Inches(3.48),Inches(11.7),Inches(2.05),LIGHT)
emails=[f"n8n{1000+i}@tertiaryinfotech.com" for i in range(1,11)]
for j,em in enumerate(emails[:5]):
    txt(s,Inches(1.1),Inches(3.60+j*0.37),Inches(5.5),Inches(0.34),[
        [(em,14,INK,False)]])
for j,em in enumerate(emails[5:]):
    txt(s,Inches(6.7),Inches(3.60+j*0.37),Inches(5.5),Inches(0.34),[
        [(em,14,INK,False)]])

# ── Password bar ────────────────────────────────────────────────
rect(s,Inches(0.85),Inches(5.70),Inches(11.7),Inches(0.65),BLUE)
txt(s,Inches(0),Inches(5.70),SW,Inches(0.65),[
    [("Password:  ",17,WHITE,True),
     ("Tertiary@888",17,WHITE,False),
     ("   ·   same for all accounts",14,LBLUE,False)]],
    align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ── Footer ──────────────────────────────────────────────────────
txt(s,Inches(0.55),Inches(7.08),Inches(6.2),Inches(0.3),[
    [("Agentic AI Automation with n8n  ·  TGS-2023035977",8.5,GREY,False)]])
txt(s,Inches(6.0),Inches(7.08),Inches(6),Inches(0.3),[
    [("© 2026 Tertiary Infotech Academy Pte Ltd",8.5,GREY,False)]],align=PP_ALIGN.CENTER)

OUT=os.path.join(REPO,"courseware","n8n-login-slide.pptx")
prs.save(OUT)
print("Saved:",OUT)
