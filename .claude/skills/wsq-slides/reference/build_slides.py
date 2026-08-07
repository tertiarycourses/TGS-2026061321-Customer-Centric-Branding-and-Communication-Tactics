#!/usr/bin/env python3
"""Generate the Customer-Centric Branding and Communication Tactics slide deck
(all-white Tertiary WSQ house style).

Design helpers are the shared wsq-slides visual component library (cover,
section, content, two_col, cards3, tile_grid, flow_h, trainer_slide,
big_statement, lead, callout, compare_table, stat_band, playbook, img_slide,
activity_slide — ONE workflow slide per activity, no step-by-step runs —
lms_slide, brk). Content is
driven entirely by course_data.py + data_domain1..4.py so the deck stays
100% aligned with the Lesson Plan and Learner Guide. Also writes
slide_map.json (topic/activity/admin-anchor -> page number) so the Lesson
Plan builder can cite the correct deck page for every teaching row.
"""
import os, sys, json, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_brandtrust import BRAND_TRUST
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4

REPO = os.path.dirname(os.path.dirname(HERE))
ASSETS = os.path.join(REPO, "courseware", "assets")

# ---------------- palette (house standard) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates (shared wsq-slides component library) ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    rect(s,Inches(10.7),Inches(0.72),Inches(1.85),Inches(1.0),BLUE)
    txt(s,Inches(10.7),Inches(0.9),Inches(1.85),Inches(0.4),[[("WSQ",22,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(10.7),Inches(1.3),Inches(1.85),Inches(0.35),[[("COURSE",9,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,38,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    """House hard rule: content slides are card-format — numbered visual tiles,
    never a plain bullet wall. (bullets() remains for two_col/cards3 panels.)"""
    plain=[(it[0] if isinstance(it,tuple) else it) for it in items]
    tsz=15 if max(len(x) for x in plain)<=140 else 13
    return tile_grid(title,plain,kicker=kicker,cols=1,size=tsz)
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    """Renders 2 or 3 evenly-spaced cards depending on how many are passed —
    never pads with an empty placeholder card."""
    s=head(slide(),title,kicker); n=min(len(cards),3)
    TOTW=Inches(11.63); X0=Inches(0.85); gap=Inches(0.33)
    cw=int((TOTW-gap*(n-1))/n)
    xs=[int(X0+(cw+gap)*i) for i in range(n)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),cw,Inches(4.7),LIGHT); rect(s,x,Inches(1.95),cw,Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),cw-Inches(0.5),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),cw-Inches(0.5),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def lead(s,text,y=1.82):
    """Copilot-style one-line intro sentence under the title rule."""
    txt(s,Inches(0.87),Inches(y),Inches(11.6),Inches(0.4),[[(text,13,GREY,False)]])
def callout(s,text,y=5.95,h=0.85,color=BLUE,size=14):
    """Full-width emphasis band with a coloured left accent."""
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(h),LIGHT)
    rect(s,Inches(0.85),Inches(y),Inches(0.1),Inches(h),color)
    txt(s,Inches(1.2),Inches(y),Inches(11.0),Inches(h),[[(text,size,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
def flow_h(title,steps,kicker=None,color=BLUE,note=None,intro=None):
    """Horizontal numbered flow: coloured chips connected by chevrons.
    Optional intro line under the title and note band under the chips."""
    s=head(slide(),title,kicker,kcolor=color)
    if intro: lead(s,intro)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n)
    y=Inches(2.45) if note else Inches(2.55)
    ch=Inches(3.0) if note else Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.38)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.38)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.45)),cw-Inches(0.32),int(ch-Inches(1.6)),[[(st,13,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note: callout(s,note,y=5.75,h=0.9,color=color)
    footer(s); return s
def compare_table(title,colheads,rows,kicker=None,intro=None):
    """Two-column comparison table with coloured column headers
    (imported design: reference-deck paradigm-shift table)."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    X0=Inches(0.85); LW=Inches(2.35); CW=Inches(4.64); Y0=Inches(2.3)
    HH=Inches(0.5); n=len(rows); rh=int((Inches(6.7)-Y0-HH)/n)
    heads=[(BLUE,colheads[0]),(TEAL,colheads[1])]
    for ci,(col,txt_) in enumerate(heads):
        x=int(X0+LW+CW*ci)
        rect(s,x,Y0,CW,HH,col)
        txt(s,x,Y0,CW,HH,[[(txt_,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    for ri,(label,a,b) in enumerate(rows):
        y=int(Y0+HH+rh*ri); bg=LIGHT if ri%2==0 else WHITE
        rect(s,X0,y,LW,rh,bg,line=LINE)
        txt(s,X0+Inches(0.15),y,LW-Inches(0.3),rh,[[(label,14,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
        for ci,cell in enumerate((a,b)):
            x=int(X0+LW+CW*ci)
            rect(s,x,y,CW,rh,bg,line=LINE)
            txt(s,x+Inches(0.18),y,CW-Inches(0.36),rh,[[(cell,13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def stat_band(title,stats,kicker=None,intro=None,note=None):
    """Big-number stat tiles + optional emphasis band (imported design:
    reference-deck economics/emotional-branding stat slides)."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    n=len(stats); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.4)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.45); ch=Inches(3.1)
    for i,(big,label) in enumerate(stats):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.12),col)
        txt(s,x,int(y+Inches(0.45)),cw,Inches(1.3),[[(big,52,col,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.25),int(y+Inches(1.85)),cw-Inches(0.5),Inches(1.1),[[(label,15,INK,False)]],align=PP_ALIGN.CENTER)
    if note: callout(s,note,y=5.85,h=0.85,color=AMBER)
    footer(s); return s
def img_slide(title,path,kicker=None,intro=None):
    """Full-width imported visual under the house header."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    p=os.path.join(ASSETS,path)
    if os.path.exists(p):
        ih=Inches(4.55); iw=int(ih*1376/768)
        s.shapes.add_picture(p,int((SW-iw)/2),Inches(2.3),height=ih)
    footer(s); return s
def playbook(title,items,kicker=None,tagline=None):
    """Numbered 01..0n columns with big faint numerals + tagline band
    (imported design: reference-deck executive playbook)."""
    s=head(slide(),title,kicker)
    n=len(items); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.0)
    for i,(h,body) in enumerate(items):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        txt(s,x,y,cw,Inches(0.9),[[("0%d"%(i+1),44,RGBColor(0xD7,0xDF,0xEA),True)]])
        txt(s,x,int(y+Inches(0.95)),cw,Inches(0.75),[[(h,16,col,True)]])
        txt(s,x,int(y+Inches(1.72)),cw,Inches(1.9),[[(body,12,INK,False)]])
    if tagline: callout(s,tagline,y=5.95,h=0.8,color=BLUE,size=16)
    footer(s); return s
def activity_slide(tag,title,desc,steps,build,duration,kicker):
    """One-slide activity brief (replaces the old overview + one-step-per-slide
    + debrief run): scenario, compact workflow strip, deliverable band."""
    if len(steps)>6:
        extra=len(steps)-5
        steps=list(steps[:5])+[f"+{extra} more steps — see the Learner Guide / lab sheet"]
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.9),Inches(0.46),TEAL)
    txt(s,Inches(0.85),Inches(1.85),Inches(1.9),Inches(0.46),[[(tag,14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(0.85),Inches(2.45),Inches(11.7),Inches(0.95),[[(desc,15,INK,False)]])
    txt(s,Inches(0.85),Inches(3.42),Inches(6),Inches(0.32),[[("WORKFLOW",12,TEAL,True)]])
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.26)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(3.78); ch=Inches(1.78); bd=Inches(0.46)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.07),TEAL)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.16)),bd,bd,TEAL)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.16)),bd,bd,[[(str(i+1),16,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.12),int(y+Inches(0.72)),cw-Inches(0.24),int(ch-Inches(0.82)),[[(st,11,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.05)),int(y+ch/2-Inches(0.25)),int(gap+Inches(0.1)),Inches(0.5),
                [[("▶",11,TEAL,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    by=Inches(5.75); bh=Inches(0.98)
    rect(s,Inches(0.85),by,Inches(11.63),bh,LIGHT); rect(s,Inches(0.85),by,Inches(0.1),bh,TEAL)
    txt(s,Inches(1.2),int(by+Inches(0.1)),Inches(8.6),Inches(0.3),[[("YOU'LL PRODUCE",11,TEAL,True)]])
    txt(s,Inches(1.2),int(by+Inches(0.4)),Inches(8.6),Inches(0.55),[[(build,13,INK,True)]])
    txt(s,Inches(9.9),by,Inches(2.35),bh,[[("Duration",11,GREY,True)],[(duration,14,INK,True)]],
        align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE,space=2)
    footer(s); return s
def lms_slide():
    """Download Course Material — visual portal card + numbered steps
    (never a bare text link)."""
    s=head(slide(),"Courseware & Assessment on the LMS","COURSE PORTAL")
    lead(s,"Everything for the open-book assessment — slides, Learner Guide and submissions — lives on the LMS.")
    lx=Inches(0.85); lw=Inches(4.1); ly=Inches(2.35); lh=Inches(3.9)
    rect(s,lx,ly,lw,lh,RGBColor(0x0F,0x2A,0x5F))
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,int(lx+Inches(0.4)),int(ly+Inches(0.35)),height=Inches(0.55))
    txt(s,lx,int(ly+Inches(1.15)),lw,Inches(0.5),[[("Tertiary Infotech Academy",18,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx,int(ly+Inches(1.7)),lw,Inches(0.4),[[("LMS and TMS for SSG courses",12,RGBColor(0xB8,0xC7,0xE0),False)]],align=PP_ALIGN.CENTER)
    rect(s,int(lx+Inches(1.05)),int(ly+Inches(2.35)),Inches(2.0),Inches(0.55),BLUE)
    txt(s,int(lx+Inches(1.05)),int(ly+Inches(2.35)),Inches(2.0),Inches(0.55),[[("Sign in / OTP",14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx,int(ly+Inches(3.15)),lw,Inches(0.4),[[("lms-tms.tertiaryinfotech.com",12,RGBColor(0x9C,0xDC,0xFE),True)]],align=PP_ALIGN.CENTER)
    items=[("Sign in","Log in with your registered email (OTP or password)."),
           ("Download materials","Slides and the Learner Guide from the course page — your open-book references."),
           ("Submit answers","Upload your Written Assessment and Case Study answers on the LMS."),
           ("TRAQOM survey","Complete the TRAQOM survey via the QR code shown on the LMS.")]
    rx=Inches(5.25); rw=Inches(7.23); gy=Inches(0.22); th=int((lh-gy*3)/4)
    for i,(h,b) in enumerate(items):
        y=int(ly+(th+gy)*i); col=PALETTE[i%len(PALETTE)]; bd=Inches(0.5)
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        oval(s,int(rx+Inches(0.25)),int(y+th/2-bd/2),bd,bd,col)
        txt(s,int(rx+Inches(0.25)),int(y+th/2-bd/2),bd,bd,[[(str(i+1),15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,int(rx+Inches(0.95)),y,int(rw-Inches(1.15)),th,[[(h,14,INK,True)],[(b,11,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    callout(s,"https://lms-tms.tertiaryinfotech.com/",y=6.42,h=0.5,color=TEAL,size=13)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value -> fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,duration,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.9),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.9),Inches(0.4),[[(tag,15,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,21,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll produce",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.9),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.75),Inches(11.2),Inches(0.5),[[("Duration:  ",13,GREY,True),(duration,13,GREY,False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(2.4),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Debrief it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.8),[[(text,17,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
SLIDE_MAP={}
cover()

# ---------------- ADMIN (front) ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="SSG DIGITAL ATTENDANCE")
SLIDE_MAP["digital_attendance_front"]=PAGE["n"]
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Certification / Credentials","20+ years of training and industry experience."),
  ("Delivers","WSQ courses on customer-centric branding, marketing communications and business skills."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with branding, marketing or communications (if any).",
 "A brand you personally admire — and why it works on you."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Actively participate — no question is too small.",
 "Respect each other's views: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required for funding eligibility."],
 kicker="HOUSEKEEPING",cols=2,size=15)
tile_grid("Skills Framework",[
 ("TSC Title",C.TSC_TITLE),
 ("TSC Code",C.TSC_CODE),
 ("Proficiency Level",C.TSC_LEVEL),
 ("Structure","4 Learning Units (LU1–LU4), each mapped to a Learning Outcome.")],
 kicker="COURSE ACCREDITATION",cols=2,size=15)
tile_grid("Course Outline",
 [(f"{t['code']} — {t['title']}", t["subtitle"]) for t in C.TOPICS],
 kicker="4 LEARNING UNITS",cols=2,size=14)
SLIDE_MAP["lesson_plan_slide"]=None
two_col(f"Lesson Plan — {C.DAYS} Days",
 [(f"Day 1 — {C.DAY_THEMES[1]}",0),
  ("LU1: Stakeholders and Organisation (Activities 1–4)",1),
  ("LU2: Customer Influence (Activities 5–8)",1),
  ("LU3: Branding in Marketing — begins (Activities 9–12)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
  ("LU3: Branding in Marketing — continued",1),
  ("LU4: Branding Effectiveness (Activities 13–17)",1),
  ("Assessment: WA (1h) + Case Study (1h), 4:00–6:00pm",1),
  ("Daily timing",0),("Day 1: 9:30am–6:30pm, 1-hour lunch. Day 2: 9:30am–6:00pm, 45-min lunch (to fit the assessment).",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2 & Assessment")
SLIDE_MAP["lesson_plan_slide"]=PAGE["n"]
tile_grid("Learning Outcomes",[
 ("Stakeholders & Organisation","Identify stakeholders/audiences and draft branding designs & reputation assessments."),
 ("Customer Influence","Capture and analyse customer insight through active listening."),
 ("Branding in Marketing","Execute branding campaigns, events and PR activities that build awareness."),
 ("Branding Effectiveness","Evaluate reputation and PR performance against KPIs and recommend improvements.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
SLIDE_MAP["briefing_assessment"]=PAGE["n"]
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
SLIDE_MAP["assessment_front"]=PAGE["n"]
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
lms_slide()

# ---------------- TOPICS + ACTIVITIES ----------------
COLOR_BY_NAME={"BLUE":BLUE,"TEAL":TEAL,"VIOLET":VIOLET,"AMBER":AMBER}
def render_insight(kind,spec):
    """Render one concept-enrichment slide imported from the reference deck
    (data_brandtrust.py) using the shared component library."""
    if kind=="pillars":
        s=head(slide(),spec["title"],spec.get("kicker"))
        if spec.get("intro"): lead(s,spec["intro"])
        items=spec["items"]; cols=2; rows=math.ceil(len(items)/cols)
        X0=Inches(0.85); Y0=Inches(2.25); TOTW=Inches(11.63); AREAH=Inches(4.45)
        gx=Inches(0.3); gy=Inches(0.26)
        cw=int((TOTW-gx*(cols-1))/cols); chh=int((AREAH-gy*(rows-1))/rows)
        for i,(h,b) in enumerate(items):
            r=i//cols; c=i%cols
            x=int(X0+(cw+gx)*c); y=int(Y0+(chh+gy)*r); col=PALETTE[i%len(PALETTE)]
            rect(s,x,y,cw,chh,LIGHT); rect(s,x,y,Inches(0.1),chh,col)
            txt(s,x+Inches(0.3),int(y+Inches(0.18)),cw-Inches(0.55),Inches(0.4),[[(h,16,col,True)]])
            txt(s,x+Inches(0.3),int(y+Inches(0.62)),cw-Inches(0.55),int(chh-Inches(0.75)),[[(b,12,INK,False)]])
        footer(s)
    elif kind=="table":
        compare_table(spec["title"],spec["colheads"],spec["rows"],
                      kicker=spec.get("kicker"),intro=spec.get("intro"))
    elif kind=="stats":
        stat_band(spec["title"],spec["stats"],kicker=spec.get("kicker"),
                  intro=spec.get("intro"),note=spec.get("note"))
    elif kind=="image":
        img_slide(spec["title"],spec["path"],kicker=spec.get("kicker"),intro=spec.get("intro"))
    elif kind=="quote":
        big_statement(spec["line1"],spec["line2"],spec["kicker"],color=VIOLET)
    elif kind=="flow":
        flow_h(spec["title"],spec["steps"],kicker=spec.get("kicker"),
               color=COLOR_BY_NAME.get(spec.get("color","BLUE"),BLUE),
               note=spec.get("note"),intro=spec.get("intro"))
    elif kind=="playbook":
        playbook(spec["title"],spec["items"],kicker=spec.get("kicker"),tagline=spec.get("tagline"))
    elif kind=="twocol":
        two_col(spec["title"],spec["left"],spec["right"],kicker=spec.get("kicker"),
                lhead=spec.get("lhead",""),rhead=spec.get("rhead",""))

TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"{t['code']}", t["title"], t["code"], t["subtitle"])
    SLIDE_MAP[f"topic{t['num']}_section"]=PAGE["n"]
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"{t['code']} · KEY CONCEPTS", cols=2, size=14)
    for kind,spec in BRAND_TRUST.get(t["num"],[]):
        render_insight(kind,spec)
    acts=TOPIC_ACTS[t["num"]]
    ngroups=2 if len(acts)<=4 else 3
    size=math.ceil(len(acts)/ngroups)
    groups=[acts[i:i+size] for i in range(0,len(acts),size)][:ngroups]
    cards=[(CARD_COLORS[gi], f"Activities {g[0]['num']}–{g[-1]['num']}", [a["title"] for a in g])
           for gi,g in enumerate(groups)]
    cards3(f"In-Class Activities — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        activity_slide(f"ACTIVITY {a['num']}", a["title"], a["desc"],
                       [st for st,_cmd in a["steps"]],
                       a["build"], a["duration"], kicker=f"{t['code']} · ACTIVITY {a['num']}")
        SLIDE_MAP[f"activity{a['num']}"]=PAGE["n"]
    recap_items=list({x["objective"]:x for x in acts}.values())[:6]
    tile_grid(f"Recap — {t['title']}",
              [(a["title"], f"You can now: {a['objective']}.") for a in recap_items],
              kicker=f"{t['code']} RECAP", cols=2, size=14)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Stakeholders & Organisation","Mapped stakeholders and audiences; drafted brand designs and reputation assessments."),
 ("Customer Influence","Captured and documented customer perception through active listening."),
 ("Branding in Marketing","Executed branding and PR campaigns aligned to strategy, plan and budget."),
 ("Branding Effectiveness","Measured reputation and PR performance against KPIs and proposed improvements.")],
 kicker="LEARNING OUTCOMES",cols=2,size=15)
content("Assessment",[
 "Written Assessment (WA, SAQ) — 1 hour.  Case Study (CS) — 1 hour.",
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to complete the TRAQOM survey and take the Assessment digital attendance (two separate steps).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="WRAP-UP")
SLIDE_MAP["assessment_end"]=PAGE["n"]
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="SSG DIGITAL ATTENDANCE")
SLIDE_MAP["digital_attendance_end"]=PAGE["n"]
big_statement("Thank You!","You are now equipped to build customer-centric brand communication that earns trust and drives results.",
              "SEE YOU AT THE NEXT ONE",color=TEAL)

# house rule: uniform fade transition on every slide (python-pptx has no API)
from pptx.oxml.ns import qn
def add_transitions(prs,kind="fade",speed="med"):
    for sl in prs.slides:
        el=sl._element
        for tr in el.findall(qn("p:transition")): el.remove(tr)
        tr=el.makeelement(qn("p:transition"),{"spd":speed})
        tr.append(el.makeelement(qn(f"p:{kind}"),{}))
        el.append(tr)
add_transitions(prs)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
with open(os.path.join(HERE,"slide_map.json"),"w") as f:
    json.dump(SLIDE_MAP,f,indent=2)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print("Saved slide_map.json")
