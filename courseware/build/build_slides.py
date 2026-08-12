#!/usr/bin/env python3
"""Generate the Customer-Centric Branding and Communication Tactics slide deck
(all-white Tertiary WSQ house style — plain flat cards, chevron flows,
comparison tables and stat bands; no native diagram library, no AI-rendered
infographics).

Design helpers are the shared wsq-slides plain component library (cover,
section, content, two_col, tile_grid, flow_h, compare_table, stat_band,
playbook, img_slide, lms_slide, trainer_slide, big_statement) plus this
course's case-study/role-play activity pattern (sub_divider, scenario_slide,
role_cards_slide, discussion_slide, reflection_slide, case_study_callout —
never step-by-step). Content is driven entirely by course_data.py +
data_domain1..4.py + data_brandtrust.py so the deck stays 100% aligned with
the Lesson Plan and Learner Guide. Also writes slide_map.json (topic/activity/
admin-anchor -> page number) so the Lesson Plan builder can cite the correct
deck page for every teaching row.
"""
import os, sys, json, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import motion

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_brandtrust import BRAND_TRUST, COMMS
ACT = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4

REPO = os.path.dirname(os.path.dirname(HERE))
ASSETS = os.path.join(REPO, "courseware", "assets")

# ---------------- palette (house standard) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
PALETTE=[BLUE,TEAL,VIOLET,AMBER]

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    size=27 if len(title)<=55 else 21
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,size,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None
def _source_line(s,text,y=6.85):
    """Attribution line. `y` is lifted by callers whose slide also carries a
    callout band at the foot, so the two never overlap."""
    if text: txt(s,Inches(0.85),Inches(y),Inches(11.6),Inches(0.3),[[(text,10,GREY,False)]])

# ---------------- slide templates (shared plain component library) ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    rect(s,Inches(10.7),Inches(0.72),Inches(1.85),Inches(1.0),BLUE)
    txt(s,Inches(10.7),Inches(0.9),Inches(1.85),Inches(0.4),[[("WSQ",22,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(10.7),Inches(1.3),Inches(1.85),Inches(0.35),[[("COURSE",9,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,38,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    """House hard rule: content slides are card-format — numbered visual tiles,
    never a plain bullet wall."""
    plain=[(it[0] if isinstance(it,tuple) else it) for it in items]
    tsz=15 if max(len(x) for x in plain)<=140 else 13
    return tile_grid(title,plain,kicker=kicker,cols=1,size=tsz)
def two_col(title,left,right,kicker=None,lhead="",rhead="",source=None):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL)
    if source: _source_line(s,source)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def lead(s,text,y=1.82):
    """Copilot-style one-line intro sentence under the title rule."""
    txt(s,Inches(0.87),Inches(y),Inches(11.6),Inches(0.4),[[(text,13,GREY,False)]])
def callout(s,text,y=5.95,h=0.85,color=BLUE,size=14):
    """Full-width emphasis band with a coloured left accent."""
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(h),LIGHT)
    rect(s,Inches(0.85),Inches(y),Inches(0.1),Inches(h),color)
    txt(s,Inches(1.2),Inches(y),Inches(11.0),Inches(h),[[(text,size,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
def flow_h(title,steps,kicker=None,color=BLUE,note=None,intro=None):
    """Horizontal numbered flow: coloured chips connected by chevrons. Each
    step is either a plain string or a (heading, body) tuple. Optional intro
    line under the title and note band under the chips."""
    s=head(slide(),title,kicker,kcolor=color)
    if intro: lead(s,intro)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.3)
    cw=int((TOTW-gap*(n-1))/n)
    y=Inches(2.45) if note else Inches(2.55)
    ch=Inches(3.0) if note else Inches(3.15); bd=Inches(0.72)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,[[(str(i+1),24,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        if isinstance(st,tuple):
            head_txt,body_txt=st
            txt(s,x+Inches(0.14),int(y+Inches(1.2)),cw-Inches(0.28),Inches(0.4),[[(head_txt,12.5,INK,True)]],align=PP_ALIGN.CENTER)
            txt(s,x+Inches(0.14),int(y+Inches(1.62)),cw-Inches(0.28),int(ch-Inches(1.75)),[[(body_txt,10.5,GREY,False)]],align=PP_ALIGN.CENTER)
        else:
            txt(s,x+Inches(0.16),int(y+Inches(1.45)),cw-Inches(0.32),int(ch-Inches(1.6)),[[(st,13,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.02)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.04)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note: callout(s,note,y=5.75,h=0.9,color=color)
    footer(s); return s
def compare_table(title,colheads,rows,kicker=None,intro=None):
    """Two-column comparison table with coloured column headers."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    X0=Inches(0.85); LW=Inches(2.35); CW=Inches(4.64); Y0=Inches(2.3)
    HH=Inches(0.5); n=len(rows); rh=int((Inches(6.7)-Y0-HH)/n)
    heads=[(BLUE,colheads[0]),(TEAL,colheads[1])]
    for ci,(col,txt_) in enumerate(heads):
        x=int(X0+LW+CW*ci)
        rect(s,x,Y0,CW,HH,col)
        txt(s,x,Y0,CW,HH,[[(txt_,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    for ri,(label,a,b) in enumerate(rows):
        y=int(Y0+HH+rh*ri); bg=LIGHT if ri%2==0 else WHITE
        rect(s,X0,y,LW,rh,bg,line=LINE)
        txt(s,X0+Inches(0.15),y,LW-Inches(0.3),rh,[[(label,14,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
        for ci,cell in enumerate((a,b)):
            x=int(X0+LW+CW*ci)
            rect(s,x,y,CW,rh,bg,line=LINE)
            txt(s,x+Inches(0.18),y,CW-Inches(0.36),rh,[[(cell,13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def stat_band(title,stats,kicker=None,intro=None,note=None):
    """Big-number stat tiles + optional emphasis band."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    n=len(stats); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.4)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.45); ch=Inches(3.1)
    for i,(big,label) in enumerate(stats):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.12),col)
        sz=44 if len(str(big))<=6 else 30
        txt(s,x,int(y+Inches(0.45)),cw,Inches(1.3),[[(str(big),sz,col,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.2),int(y+Inches(1.85)),cw-Inches(0.4),Inches(1.1),[[(label,13,INK,False)]],align=PP_ALIGN.CENTER)
    if note: callout(s,note,y=5.85,h=0.85,color=AMBER,size=12.5)
    footer(s); return s
def img_slide(title,path,kicker=None,intro=None):
    """Full-width imported visual under the house header."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    p=os.path.join(ASSETS,path)
    if os.path.exists(p):
        ih=Inches(4.55); iw=int(ih*1376/768)
        s.shapes.add_picture(p,int((SW-iw)/2),Inches(2.3),height=ih)
    footer(s); return s

# ---------------- native vector diagrams (replace imported raster art) --------
# These three replace the dark AI-rendered PNGs that clashed with the all-white
# house theme. Drawn with the same rect/oval/txt primitives as every other
# slide, so they stay on-palette, stay crisp at any zoom, and stay editable.

def funnel_vs_journey(title,kicker=None,intro=None,note=None):
    """Side-by-side: the linear marketing funnel vs the looping customer journey."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    Y0=Inches(2.35); PANH=Inches(3.75); PW=Inches(5.68); gap=Inches(0.27)
    LX=Inches(0.85); RX=int(LX+PW+gap)

    # ---- left panel: the funnel (business-centric)
    rect(s,LX,Y0,PW,PANH,LIGHT); rect(s,LX,Y0,PW,Inches(0.1),BLUE)
    txt(s,LX,int(Y0+Inches(0.22)),PW,Inches(0.34),
        [[("THE MARKETING FUNNEL",13,BLUE,True)]],align=PP_ALIGN.CENTER)
    txt(s,LX,int(Y0+Inches(0.58)),PW,Inches(0.3),
        [[("Business-centric  ·  linear  ·  ends at the sale",10.5,GREY,False)]],align=PP_ALIGN.CENTER)
    stages=["AWARENESS","INTEREST","CONSIDERATION","INTENT","PURCHASE"]
    top=int(Y0+Inches(1.0)); bh=Inches(0.4); bgap=Inches(0.07)
    wide=Inches(4.5); narrow=Inches(1.9)
    for i,st in enumerate(stages):
        frac=i/(len(stages)-1)
        w=int(wide-(wide-narrow)*frac)
        x=int(LX+PW/2-w/2); y=int(top+(bh+bgap)*i)
        # tint deepens toward the point of the funnel
        col=BLUE if i==len(stages)-1 else RGBColor(0x8A,0xB4,0xF0)
        rect(s,x,y,w,bh,col)
        txt(s,x,y,w,bh,[[(st,10.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,LX+Inches(0.3),int(Y0+PANH-Inches(0.5)),PW-Inches(0.6),Inches(0.42),
        [[("Tracks drop-off and acquisition cost.",10.5,INK,False)]],align=PP_ALIGN.CENTER)

    # ---- right panel: the journey loop (customer-centric)
    rect(s,RX,Y0,PW,PANH,LIGHT); rect(s,RX,Y0,PW,Inches(0.1),TEAL)
    txt(s,RX,int(Y0+Inches(0.22)),PW,Inches(0.34),
        [[("THE CUSTOMER JOURNEY",13,TEAL,True)]],align=PP_ALIGN.CENTER)
    txt(s,RX,int(Y0+Inches(0.58)),PW,Inches(0.3),
        [[("Customer-centric  ·  non-linear  ·  loops into advocacy",10.5,GREY,False)]],align=PP_ALIGN.CENTER)
    # ring of stage chips — a loop, drawn as an ellipse of nodes
    ring=["DISCOVERY","RESEARCH","PURCHASE","EXPERIENCE","RETENTION","ADVOCACY"]
    cx=int(RX+PW/2); cy=int(Y0+Inches(2.05))
    rx=Inches(2.02); ry=Inches(0.9)
    cw_=Inches(1.42); ch_=Inches(0.42)
    for i,st in enumerate(ring):
        ang=-math.pi/2+2*math.pi*i/len(ring)
        x=int(cx+rx*math.cos(ang)-cw_/2); y=int(cy+ry*math.sin(ang)-ch_/2)
        col=TEAL if st in ("ADVOCACY","DISCOVERY") else RGBColor(0x5A,0xCf,0xA8)
        rect(s,x,y,cw_,ch_,col)
        txt(s,x,y,cw_,ch_,[[(st,9,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # centre hub makes the loop legible as a cycle
    hub=Inches(0.98)
    oval(s,int(cx-hub/2),int(cy-hub/2),hub,hub,WHITE)
    txt(s,int(cx-hub/2),int(cy-hub/2),hub,hub,
        [[("↻",20,TEAL,True)],[("loops",8.5,GREY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
    txt(s,RX+Inches(0.3),int(Y0+PANH-Inches(0.5)),PW-Inches(0.6),Inches(0.42),
        [[("Advocacy feeds the next customer's discovery.",10.5,INK,False)]],align=PP_ALIGN.CENTER)

    callout(s,note or "The funnel explains where customers drop off. The journey explains why.",
            y=6.35,h=0.62,color=VIOLET,size=13)
    footer(s); return s


def journey_loop(title,kicker=None,intro=None,stages=None,note=None):
    """The 5-stage journey as a horizontal chevron that visibly returns to the start."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    stages=stages or [
        ("Awareness","The customer first encounters the brand."),
        ("Consideration","They compare it against alternatives."),
        ("Decision","They buy — the funnel would stop here."),
        ("Retention","The experience decides whether they stay."),
        ("Advocacy","They tell others, feeding new awareness."),
    ]
    n=len(stages); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.26)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(2.65); bd=Inches(0.66)
    for i,(h_,b_) in enumerate(stages):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),col)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.32)),bd,bd,col)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.32)),bd,bd,
            [[(str(i+1),21,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.12),int(y+Inches(1.12)),cw-Inches(0.24),Inches(0.38),
            [[(h_,12.5,INK,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.12),int(y+Inches(1.52)),cw-Inches(0.24),int(ch-Inches(1.64)),
            [[(b_,10,GREY,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.02)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.04)),Inches(0.6),
                [[("▶",14,PALETTE[i%len(PALETTE)],True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # the return arc — what makes it a loop rather than a funnel
    arcy=int(y+ch+Inches(0.18))
    rect(s,X0,arcy,TOTW,Inches(0.06),TEAL)
    rect(s,X0,arcy,Inches(0.06),Inches(0.3),TEAL)
    rect(s,int(X0+TOTW-Inches(0.06)),arcy,Inches(0.06),Inches(0.3),TEAL)
    txt(s,X0,int(arcy+Inches(0.12)),TOTW,Inches(0.42),
        [[("◀  advocacy feeds awareness — the journey loops, it does not end",11.5,TEAL,True)]],
        align=PP_ALIGN.CENTER)
    if note: callout(s,note,y=6.42,h=0.58,color=BLUE,size=12.5)
    footer(s); return s


def trust_architecture(title,kicker=None,intro=None,layers=None,note=None):
    """Concentric model: the customer at the centre, ringed by the three
    layers that build modern brand trust."""
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    layers=layers or [
        ("Emotional Authenticity","Values, storytelling, transparency",VIOLET),
        ("The Customer Journey Loop","Awareness through to advocacy",TEAL),
        ("AI Brand Authority","Identity, access, governance",BLUE),
    ]
    cx=int(SW/2); cy=Inches(4.32)
    # rings drawn largest-first so inner rings paint on top
    sizes=[Inches(3.95),Inches(3.0),Inches(2.05)]
    for (lbl,desc,col),d in zip(reversed(layers),sizes):
        oval(s,int(cx-d/2),int(cy-d/2),d,d,col)
    core=Inches(1.24)
    oval(s,int(cx-core/2),int(cy-core/2),core,core,WHITE)
    txt(s,int(cx-core/2),int(cy-core/2),core,core,
        [[("THE",9,GREY,True)],[("CUSTOMER",11,INK,True)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)

    # legend cards to the left and right — keeps labels off the rings
    lw=Inches(3.55); lh=Inches(1.16)
    spots=[(Inches(0.85),Inches(2.62)),(Inches(0.85),Inches(5.12)),
           (int(SW-Inches(0.85)-lw),Inches(3.87))]
    for (lbl,desc,col),(lx,ly) in zip(layers,spots):
        rect(s,lx,ly,lw,lh,LIGHT); rect(s,lx,ly,Inches(0.09),lh,col)
        txt(s,lx+Inches(0.26),int(ly+Inches(0.16)),lw-Inches(0.42),Inches(0.42),
            [[(lbl,13,col,True)]])
        txt(s,lx+Inches(0.26),int(ly+Inches(0.58)),lw-Inches(0.42),Inches(0.5),
            [[(desc,10.5,GREY,False)]])
    callout(s,note or "A self-sustaining ecosystem: the customer at the centre, anchored by "
                      "emotion, navigated by journey mapping, scaled within AI authority boundaries.",
            y=6.42,h=0.6,color=AMBER,size=12)
    footer(s); return s
# ---------------- native PPT charts, scorecards and process maps -------------
# Real chart parts (editable in PowerPoint: right-click > Edit Data), plus a
# KPI scorecard and a swimlane process map. All on the house palette.

def _style_chart(chart,legend=True,font=10):
    """House styling for a native chart part."""
    chart.has_title=False
    chart.font.size=Pt(font); chart.font.name="Arial"; chart.font.color.rgb=GREY
    if legend:
        chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout=False
        chart.legend.font.size=Pt(font); chart.legend.font.color.rgb=INK
    else:
        chart.has_legend=False
    return chart

def chart_slide(title,kind,categories,series,kicker=None,intro=None,note=None,
                source=None,legend=True,colors=None,number_format=None):
    """A native PowerPoint chart under the house header.

    kind       -- "bar" | "column" | "pie" | "doughnut" | "line"
    categories -- list of category labels
    series     -- list of (name, [values]) tuples
    """
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    cd=CategoryChartData(); cd.categories=categories
    for nm,vals in series: cd.add_series(nm,vals)
    ctype={"bar":XL_CHART_TYPE.BAR_CLUSTERED,
           "column":XL_CHART_TYPE.COLUMN_CLUSTERED,
           "pie":XL_CHART_TYPE.PIE,
           "doughnut":XL_CHART_TYPE.DOUGHNUT,
           "line":XL_CHART_TYPE.LINE_MARKERS}[kind]
    y=Inches(2.3); h=Inches(4.1) if (note or source) else Inches(4.45)
    gf=s.shapes.add_chart(ctype,Inches(1.4),y,Inches(10.5),h,cd)
    ch=gf.chart; _style_chart(ch,legend=legend)

    pal=colors or PALETTE
    if kind in ("pie","doughnut"):
        pts=ch.plots[0]
        pts.has_data_labels=True
        dl=pts.data_labels; dl.number_format='0"%"'; dl.number_format_is_linked=False
        dl.font.size=Pt(11); dl.font.bold=True; dl.font.color.rgb=WHITE
        for i,pt in enumerate(ch.series[0].points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb=pal[i%len(pal)]
            pt.format.line.color.rgb=WHITE; pt.format.line.width=Pt(1.5)
    else:
        for i,sr in enumerate(ch.series):
            sr.format.fill.solid(); sr.format.fill.fore_color.rgb=pal[i%len(pal)]
            sr.format.line.fill.background()
        if len(series)==1:
            ch.plots[0].has_data_labels=True
            dl=ch.plots[0].data_labels
            dl.font.size=Pt(10); dl.font.bold=True; dl.font.color.rgb=INK
            if number_format:
                dl.number_format=number_format; dl.number_format_is_linked=False
        try:
            ch.value_axis.has_major_gridlines=True
            ch.value_axis.format.line.color.rgb=LINE
            ch.category_axis.format.line.color.rgb=LINE
        except Exception:
            pass
    if note:
        callout(s,note,y=6.42,h=0.6,color=AMBER,size=12)
        _source_line(s,source,y=6.03)
    else:
        _source_line(s,source)
    footer(s); return s

def scorecard(title,cards,kicker=None,intro=None,note=None,source=None,cols=None):
    """KPI scorecard: metric tiles with value, label, target and trend chip.

    cards -- list of dicts: {value, label, target (optional), trend (optional),
             direction: "up"|"down"|"flat"}
    """
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    n=len(cards); cols=cols or (len(cards) if len(cards)<=4 else 3)
    rows=math.ceil(n/cols)
    X0=Inches(0.85); TOTW=Inches(11.63); gx=Inches(0.28); gy=Inches(0.26)
    Y0=Inches(2.42); AREAH=Inches(3.5) if (note and source) else (
                             Inches(3.85) if (note or source) else Inches(4.2))
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    ARROW={"up":"▲","down":"▼","flat":"■"}
    for i,c in enumerate(cards):
        r=i//cols; cc=i%cols
        x=int(X0+(cw+gx)*cc); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),col)
        val=str(c["value"]); vs=40 if len(val)<=5 else 30
        txt(s,x,int(y+Inches(0.3)),cw,Inches(1.0),[[(val,vs,col,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.16),int(y+Inches(1.32)),cw-Inches(0.32),Inches(0.78),
            [[(c["label"],12,INK,False)]],align=PP_ALIGN.CENTER)
        by=int(y+ch-Inches(0.62))
        if c.get("target"):
            rect(s,x+Inches(0.16),by,cw-Inches(0.32),Inches(0.44),WHITE)
            trend=c.get("trend"); d=c.get("direction","flat")
            # A leading indicator warns early; a lagging one only confirms. Colour
            # by that meaning rather than by the arrow, so the chip cannot imply
            # "good/bad" where none is meant.
            if trend=="leading":   tcol,mark=BLUE,"◆"
            elif trend=="lagging": tcol,mark=GREY,"○"
            else:                  tcol,mark=((TEAL if d=="up" else
                                               RGBColor(0xD9,0x3A,0x3A) if d=="down" else GREY),
                                              ARROW[d])
            line=[("Target: ",9.5,GREY,False),(str(c["target"]),9.5,INK,True)]
            if trend: line+=[("   "+mark+" ",9.5,tcol,True),(str(trend),9.5,tcol,True)]
            txt(s,x+Inches(0.16),by,cw-Inches(0.32),Inches(0.44),[line],
                align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note:
        callout(s,note,y=6.45,h=0.58,color=BLUE,size=12)
        _source_line(s,source,y=6.06)
    else:
        _source_line(s,source)
    footer(s); return s

def process_map(title,lanes,kicker=None,intro=None,note=None,source=None):
    """Swimlane process map: each lane is an actor, each step a chip in sequence.

    lanes -- list of (lane_name, [step, ...]); steps align to columns so the
             hand-offs between actors read left to right.
    """
    s=head(slide(),title,kicker)
    if intro: lead(s,intro)
    ncol=max(len(st) for _,st in lanes)
    X0=Inches(0.85); LANEW=Inches(2.05); TOTW=Inches(11.63)
    stepw=int((TOTW-LANEW-Inches(0.2)*(ncol-1))/ncol); gx=Inches(0.2)
    Y0=Inches(2.4); AREAH=Inches(3.85) if (note or source) else Inches(4.2)
    lh=int((AREAH-Inches(0.14)*(len(lanes)-1))/len(lanes)); gy=Inches(0.14)
    # index the occupied cells first so hand-offs between lanes can be drawn
    occupied={}   # column -> (lane index, colour)
    for li,(lane,steps) in enumerate(lanes):
        for si,st in enumerate(steps):
            if st: occupied[si]=(li,PALETTE[li%len(PALETTE)])

    def cell_x(si): return int(X0+LANEW+Inches(0.2)+(stepw+gx)*si)
    def lane_y(li): return int(Y0+(lh+gy)*li)

    for li,(lane,steps) in enumerate(lanes):
        y=lane_y(li); col=PALETTE[li%len(PALETTE)]
        rect(s,X0,y,LANEW,lh,col)
        txt(s,X0+Inches(0.14),y,LANEW-Inches(0.28),lh,[[(lane,12,WHITE,True)]],
            anchor=MSO_ANCHOR.MIDDLE)
        for si,st in enumerate(steps):
            if not st: continue
            x=cell_x(si)
            rect(s,x,y,stepw,lh,LIGHT,line=LINE); rect(s,x,y,Inches(0.07),lh,col)
            txt(s,x+Inches(0.18),y,stepw-Inches(0.3),lh,[[(st,10.5,INK,False)]],
                anchor=MSO_ANCHOR.MIDDLE)

    # connectors: same lane -> a chevron in the gap; lane change -> an elbow that
    # makes the hand-off between actors explicit rather than implied by a blank.
    cols=sorted(occupied)
    for a_,b_ in zip(cols,cols[1:]):
        la,ca=occupied[a_]; lb,_=occupied[b_]
        ya,yb=lane_y(la),lane_y(lb)
        xa_end=cell_x(a_)+stepw; xb=cell_x(b_)
        if la==lb and b_==a_+1:
            txt(s,xa_end,int(ya+lh/2-Inches(0.16)),gx,Inches(0.32),
                [[("▶",11,ca,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        else:
            mid=int(xa_end+gx/2)
            rect(s,mid,int(ya+lh/2),Inches(0.035),int(yb+lh/2-(ya+lh/2)),ca)
            rect(s,mid,int(yb+lh/2),int(xb-mid),Inches(0.035),ca)
            txt(s,int(xb-Inches(0.24)),int(yb+lh/2-Inches(0.16)),Inches(0.24),Inches(0.32),
                [[("▶",11,ca,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if note: callout(s,note,y=6.45,h=0.58,color=TEAL,size=12)
    _source_line(s,source)
    footer(s); return s

def playbook(title,items,kicker=None,tagline=None):
    """Numbered 01..0n columns with big faint numerals + tagline band."""
    s=head(slide(),title,kicker)
    n=len(items); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.0)
    for i,(h,body) in enumerate(items):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        txt(s,x,y,cw,Inches(0.9),[[("0%d"%(i+1),44,RGBColor(0xD7,0xDF,0xEA),True)]])
        txt(s,x,int(y+Inches(0.95)),cw,Inches(0.75),[[(h,16,col,True)]])
        txt(s,x,int(y+Inches(1.72)),cw,Inches(1.9),[[(body,12,INK,False)]])
    if tagline: callout(s,tagline,y=5.95,h=0.8,color=BLUE,size=16)
    footer(s); return s
def lms_slide():
    """Download Course Material — visual portal card + numbered steps."""
    s=head(slide(),"Courseware & Assessment on the LMS","COURSE PORTAL")
    lead(s,"Everything for the open-book assessment — slides, Learner Guide and submissions — lives on the LMS.")
    lx=Inches(0.85); lw=Inches(4.1); ly=Inches(2.35); lh=Inches(3.9)
    rect(s,lx,ly,lw,lh,RGBColor(0x0F,0x2A,0x5F))
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,int(lx+Inches(0.4)),int(ly+Inches(0.35)),height=Inches(0.55))
    txt(s,lx,int(ly+Inches(1.15)),lw,Inches(0.5),[[("Tertiary Infotech Academy",18,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx,int(ly+Inches(1.7)),lw,Inches(0.4),[[("LMS and TMS for SSG courses",12,RGBColor(0xB8,0xC7,0xE0),False)]],align=PP_ALIGN.CENTER)
    rect(s,int(lx+Inches(1.05)),int(ly+Inches(2.35)),Inches(2.0),Inches(0.55),BLUE)
    txt(s,int(lx+Inches(1.05)),int(ly+Inches(2.35)),Inches(2.0),Inches(0.55),[[("Sign in / OTP",14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx,int(ly+Inches(3.15)),lw,Inches(0.4),[[("lms-tms.tertiaryinfotech.com",12,RGBColor(0x9C,0xDC,0xFE),True)]],align=PP_ALIGN.CENTER)
    items=[("Sign in","Log in with your registered email (OTP or password)."),
           ("Download materials","Slides and the Learner Guide from the course page — your open-book references."),
           ("Submit answers","Upload your Written Assessment and Case Study answers on the LMS."),
           ("TRAQOM survey","Complete the TRAQOM survey via the QR code shown on the LMS.")]
    rx=Inches(5.25); rw=Inches(7.23); gy=Inches(0.22); th=int((lh-gy*3)/4)
    for i,(h,b) in enumerate(items):
        y=int(ly+(th+gy)*i); col=PALETTE[i%len(PALETTE)]; bd=Inches(0.5)
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        oval(s,int(rx+Inches(0.25)),int(y+th/2-bd/2),bd,bd,col)
        txt(s,int(rx+Inches(0.25)),int(y+th/2-bd/2),bd,bd,[[(str(i+1),15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,int(rx+Inches(0.95)),y,int(rw-Inches(1.15)),th,[[(h,14,INK,True)],[(b,11,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    callout(s,"https://lms-tms.tertiaryinfotech.com/",y=6.42,h=0.5,color=TEAL,size=13)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value -> fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s

def sub_divider(kicker,title,n=""):
    """Full-bleed light divider for a K/A sub-topic (T1, T2, ...) within a
    Learning Unit — mirrors section()'s layout (ghost watermark, coloured
    left bar) in teal so it reads as one level below the LU section divider."""
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,TEAL)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),BLUE)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,16,TEAL,True)]])
    size=34 if len(title)<=50 else (26 if len(title)<=80 else 21)
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,size,INK,True)]])
    if n: txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s); return s

# --------- case-study / role-play activity pattern (never step-by-step) ---------
CASE_TAG={"case_study":"CASE STUDY","role_play":"ROLE PLAY"}
def case_study_callout(company,headline,insight,kicker=None,color=BLUE):
    """Light 'Real-World Example' callout card citing a named company — for
    frequent use inline within a topic."""
    s=head(slide(),headline,kicker,kcolor=color)
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(4.6),LIGHT)
    rect(s,Inches(0.85),Inches(1.95),Inches(0.14),Inches(4.6),color)
    rect(s,Inches(1.2),Inches(2.3),Inches(3.1),Inches(0.55),color)
    txt(s,Inches(1.2),Inches(2.3),Inches(3.1),Inches(0.55),[[("REAL-WORLD EXAMPLE",13,WHITE,True)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(1.2),Inches(3.15),Inches(10.9),Inches(0.7),[[(company,26,INK,True)]])
    txt(s,Inches(1.2),Inches(3.95),Inches(10.9),Inches(2.4),[[(insight,16,GREY,False)]])
    footer(s); return s
def scenario_slide(kicker,act_title,case_type,narrative,build,duration):
    """The case narrative as a styled scenario card — tagged Case Study or Role
    Play — never a numbered instruction list."""
    s=head(slide(),act_title,kicker,kcolor=TEAL)
    tag=CASE_TAG.get(case_type,"CASE STUDY")
    rect(s,Inches(0.85),Inches(1.85),Inches(1.85),Inches(0.46),TEAL)
    txt(s,Inches(0.85),Inches(1.89),Inches(1.85),Inches(0.4),[[(tag,13,WHITE,True)]],align=PP_ALIGN.CENTER)
    rect(s,Inches(0.85),Inches(2.5),Inches(11.63),Inches(3.15),LIGHT); rect(s,Inches(0.85),Inches(2.5),Inches(0.1),Inches(3.15),TEAL)
    txt(s,Inches(1.2),Inches(2.72),Inches(10.9),Inches(0.35),[[("THE SITUATION",12,TEAL,True)]])
    paras=[[(p,14.5,INK,False)] for p in narrative]
    txt(s,Inches(1.2),Inches(3.1),Inches(11.0),Inches(2.45),paras,space=8)
    rect(s,Inches(0.85),Inches(5.85),Inches(11.63),Inches(0.75),WHITE,line=LINE)
    txt(s,Inches(1.1),Inches(5.85),Inches(8.0),Inches(0.75),
        [[("You'll produce:  ",13,BLUE,True),(build,13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(9.6),Inches(5.85),Inches(2.85),Inches(0.75),
        [[("Duration: ",13,GREY,True),(duration,13,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
    footer(s); return s
def role_cards_slide(kicker,act_title,roles):
    """Persona cards for role-play activities: role tag, goal, brief."""
    s=head(slide(),act_title,kicker,kcolor=TEAL)
    txt(s,Inches(0.85),Inches(1.88),Inches(9),Inches(0.4),[[("ASSIGN THE ROLES",13,GREY,True)]])
    n=len(roles); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.3)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.4); ch=Inches(4.25)
    for i,(name,goal,brief) in enumerate(roles):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.5),col)
        txt(s,x,y,cw,Inches(0.5),[[(name,15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,int(x+Inches(0.22)),int(y+Inches(0.68)),cw-Inches(0.44),Inches(0.6),[[("GOAL",11,col,True)]])
        txt(s,int(x+Inches(0.22)),int(y+Inches(0.98)),cw-Inches(0.44),Inches(0.9),[[(goal,12.5,INK,True)]])
        txt(s,int(x+Inches(0.22)),int(y+Inches(1.95)),cw-Inches(0.44),Inches(2.2),[[(brief,12,GREY,False)]])
    footer(s); return s
def discussion_slide(kicker,act_title,prompts):
    """Numbered discussion/decision prompts — open questions the group works
    through together, never literal click-by-click steps."""
    s=head(slide(),act_title,kicker,kcolor=TEAL)
    txt(s,Inches(0.85),Inches(1.88),Inches(9),Inches(0.4),[[("DISCUSSION & DECISION PROMPTS",13,GREY,True)]])
    n=len(prompts); X0=Inches(0.85); Y0=Inches(2.35); TOTW=Inches(11.63); gap=Inches(0.16)
    rowh=int((Inches(4.3)-gap*(n-1))/n); bd=min(Inches(0.55),int(rowh*0.68))
    for i,p in enumerate(prompts):
        y=int(Y0+(rowh+gap)*i)
        rect(s,X0,y,TOTW,rowh,LIGHT); rect(s,X0,y,Inches(0.08),rowh,TEAL)
        oval(s,int(X0+Inches(0.22)),int(y+rowh/2-bd/2),bd,bd,TEAL)
        txt(s,int(X0+Inches(0.22)),int(y+rowh/2-bd/2),bd,bd,[[(str(i+1),15,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=int(X0+Inches(0.22)+bd+Inches(0.28)); tw=int(TOTW-(bd+Inches(0.9)))
        txt(s,tx,y,tw,rowh,[[(p,14.5 if rowh>Inches(0.7) else 13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def reflection_slide(kicker,act_title,reflection_points,debrief_check):
    """Closing reflection questions + the facilitator's debrief check."""
    s=head(slide(),act_title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(2.9),LIGHT); rect(s,Inches(0.85),Inches(1.95),Inches(0.1),Inches(2.9),VIOLET)
    txt(s,Inches(1.15),Inches(2.15),Inches(11),Inches(0.4),[[("REFLECT & DISCUSS",13,VIOLET,True)]])
    bullets(s,Inches(1.15),Inches(2.6),Inches(11.0),Inches(2.1),reflection_points,size=14.5,mcolor=VIOLET,gap=9)
    rect(s,Inches(0.85),Inches(5.1),Inches(11.63),Inches(1.5),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.15),Inches(5.32),Inches(11),Inches(0.4),[[("✓  Facilitator debrief check",15,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.15),Inches(5.75),Inches(11.0),Inches(0.8),[[(debrief_check,13,INK,False)]])
    footer(s); return s

# --------- concept-enrichment insight slides (data_brandtrust.py) ---------
COLOR_BY_NAME={"BLUE":BLUE,"TEAL":TEAL,"VIOLET":VIOLET,"AMBER":AMBER}
def render_insight(kind,spec):
    if kind=="pillars":
        s=head(slide(),spec["title"],spec.get("kicker"))
        if spec.get("intro"): lead(s,spec["intro"])
        items=spec["items"]; cols=2; rows=math.ceil(len(items)/cols)
        X0=Inches(0.85); Y0=Inches(2.25); TOTW=Inches(11.63); AREAH=Inches(4.45)
        gx=Inches(0.3); gy=Inches(0.26)
        cw=int((TOTW-gx*(cols-1))/cols); chh=int((AREAH-gy*(rows-1))/rows)
        for i,(h,b) in enumerate(items):
            r=i//cols; c=i%cols
            x=int(X0+(cw+gx)*c); y=int(Y0+(chh+gy)*r); col=PALETTE[i%len(PALETTE)]
            rect(s,x,y,cw,chh,LIGHT); rect(s,x,y,Inches(0.1),chh,col)
            txt(s,x+Inches(0.3),int(y+Inches(0.18)),cw-Inches(0.55),Inches(0.4),[[(h,16,col,True)]])
            txt(s,x+Inches(0.3),int(y+Inches(0.62)),cw-Inches(0.55),int(chh-Inches(0.75)),[[(b,12,INK,False)]])
        footer(s)
    elif kind=="table":
        compare_table(spec["title"],spec["colheads"],spec["rows"],
                      kicker=spec.get("kicker"),intro=spec.get("intro"))
    elif kind=="stats":
        stat_band(spec["title"],spec["stats"],kicker=spec.get("kicker"),
                  intro=spec.get("intro"),note=spec.get("note"))
    elif kind=="image":
        img_slide(spec["title"],spec["path"],kicker=spec.get("kicker"),intro=spec.get("intro"))
    elif kind=="funnel_journey":
        funnel_vs_journey(spec["title"],kicker=spec.get("kicker"),
                          intro=spec.get("intro"),note=spec.get("note"))
    elif kind=="journey_loop":
        journey_loop(spec["title"],kicker=spec.get("kicker"),intro=spec.get("intro"),
                     stages=spec.get("stages"),note=spec.get("note"))
    elif kind=="architecture":
        trust_architecture(spec["title"],kicker=spec.get("kicker"),intro=spec.get("intro"),
                           layers=spec.get("layers"),note=spec.get("note"))
    elif kind=="chart":
        chart_slide(spec["title"],spec.get("chart","column"),spec["categories"],spec["series"],
                    kicker=spec.get("kicker"),intro=spec.get("intro"),note=spec.get("note"),
                    source=spec.get("source"),legend=spec.get("legend",True),
                    number_format=spec.get("number_format"))
    elif kind=="scorecard":
        scorecard(spec["title"],spec["cards"],kicker=spec.get("kicker"),intro=spec.get("intro"),
                  note=spec.get("note"),source=spec.get("source"),cols=spec.get("cols"))
    elif kind=="process_map":
        process_map(spec["title"],spec["lanes"],kicker=spec.get("kicker"),intro=spec.get("intro"),
                    note=spec.get("note"),source=spec.get("source"))
    elif kind=="quote":
        big_statement(spec["line1"],spec["line2"],spec["kicker"],color=VIOLET)
    elif kind=="flow":
        flow_h(spec["title"],spec["steps"],kicker=spec.get("kicker"),
               color=COLOR_BY_NAME.get(spec.get("color","BLUE"),BLUE),
               note=spec.get("note"),intro=spec.get("intro"))
    elif kind=="playbook":
        playbook(spec["title"],spec["items"],kicker=spec.get("kicker"),tagline=spec.get("tagline"))
    elif kind=="twocol":
        two_col(spec["title"],spec["left"],spec["right"],kicker=spec.get("kicker"),
                lhead=spec.get("lhead",""),rhead=spec.get("rhead",""))

# --------- plain theory-slot renderers (What is / How it Works / Compare / Supporting Data) ---------
def _norm_pair_items(items):
    """Normalise a what_is/visual item list to (heading, body) tuples,
    tolerating 2-tuples (h,b) or 3-tuples (icon,h,b)."""
    out=[]
    for it in items:
        if len(it)==3: out.append((it[1],it[2]))
        else: out.append((it[0],it[1]))
    return out

def render_what_is(a,t_tag):
    """'What is X?' slot — plain tile grid of the underlying concept items,
    regardless of what_is_kind (the old diagram-kind hint is unused now)."""
    title=f"What is {a['t_statement']}?"; kicker=t_tag
    if a.get("what_is_kind")=="ecosystem":
        items=[label for _icon,label in a["what_is_items"]]
    else:
        items=_norm_pair_items(a["what_is_items"])
    s=tile_grid(title,items,kicker=kicker,cols=2,size=14)
    _source_line(s,a.get("what_is_source"))
def render_process(a,t_tag):
    """HOW IT WORKS slot — plain horizontal chevron flow of the process steps."""
    kicker=f"{t_tag} · HOW IT WORKS"
    flow_h(a["process_title"],a["process_items"],kicker=kicker,color=TEAL)
def render_compare(a,t_tag):
    """COMPARE slot — two-column comparison (default) or a plain tile grid for
    the rare quadrant-shaped entries."""
    kicker=f"{t_tag} · COMPARE"
    left=a.get("compare_left"); right=a.get("compare_right")
    if isinstance(left,tuple) and len(left)==3:
        # (letter, heading, [bullet items]) shape — flatten to plain bullet lists
        left=left[2]; right=right[2]
    elif isinstance(left,list) and left and isinstance(left[0],tuple) and len(left[0])==3:
        # list of (icon, heading, body) triples — flatten each to one bullet line
        left=[f"{h} — {b}" for _icon,h,b in left]
        right=[f"{h} — {b}" for _icon,h,b in right]
    if left is not None:
        s=two_col(a["compare_title"],left,right,kicker=kicker,
                lhead=a.get("compare_lhead",""),rhead=a.get("compare_rhead",""))
        _source_line(s,a.get("compare_source"))
    elif a.get("compare_quadrants"):
        items=[(name,"; ".join(sub_items)) for name,_tag,sub_items in a["compare_quadrants"]]
        s=tile_grid(a["compare_title"],items,kicker=kicker,cols=2,size=13)
        _source_line(s,a.get("compare_source"))
    else:
        s=tile_grid(a["compare_title"],a.get("compare_items",[]),kicker=kicker,cols=2,size=14)
        _source_line(s,a.get("compare_source"))
def render_visual(a,t_tag):
    """SUPPORTING DATA slot — plain stat tiles for numeric (label, number)
    data, plain tile grids for everything else (the old bespoke-diagram kind
    hint is only consulted for the handful of kinds with a genuinely
    different data shape — kpi/funnel/journey/persona/empathy/blueprint/tree)."""
    kind=a.get("visual_kind","tile"); kicker=f"{t_tag} · SUPPORTING DATA"
    title=a["visual_title"]; src=a.get("visual_source")
    if kind=="kpi":
        stats=[(big,f"{lbl} — {cap}") for big,lbl,cap,_trend in a["visual_items"]]
        stat_band(title,stats,kicker=kicker,note=src)
    elif kind=="funnel":
        stats=[(pct,f"{stg} — {desc}") for stg,pct,desc in a["visual_items"]]
        stat_band(title,stats,kicker=kicker,note=src)
    elif kind=="journey":
        items=[(stage,desc) for _icon,stage,desc,_score in a["visual_items"]]
        s=tile_grid(title,items,kicker=kicker,cols=2,size=13); _source_line(s,src)
    elif kind=="persona":
        items=[(name,f"{role} — Wants: {'; '.join(wants)}. Concerns: {'; '.join(concerns)}")
               for name,role,wants,concerns in a["visual_items"]]
        s=tile_grid(title,items,kicker=kicker,cols=1,size=13); _source_line(s,src)
    elif kind=="empathy":
        items=[("Says","; ".join(a["visual_says"])),("Thinks","; ".join(a["visual_thinks"])),
               ("Does","; ".join(a["visual_does"])),("Feels","; ".join(a["visual_feels"]))]
        s=tile_grid(title,items,kicker=kicker,cols=2,size=13); _source_line(s,src)
    elif kind=="blueprint":
        labels=a["visual_stage_labels"]
        items=[(lane,"; ".join(f"{labels[i]}: {si}" for i,si in enumerate(stage_items)))
               for lane,_color,stage_items in a["visual_items"]]
        s=tile_grid(title,items,kicker=kicker,cols=2,size=12); _source_line(s,src)
    elif kind=="tree":
        items=[("The Fork",a["visual_root"])]+[(f"{opt} — {hd}",desc) for opt,hd,desc,_c in a["visual_items"]]
        s=tile_grid(title,items,kicker=kicker,cols=1,size=13); _source_line(s,src)
    else:
        raw=a["visual_items"]
        if len(raw[0])==2 and isinstance(raw[0][1],(int,float)):
            stats=[(f"{v}%",lbl) for lbl,v in raw]
            stat_band(title,stats,kicker=kicker,note=src)
        else:
            items=_norm_pair_items(raw)
            s=tile_grid(title,items,kicker=kicker,cols=2,size=14); _source_line(s,src)

# ============================================================ BUILD
SLIDE_MAP={}
cover()

# ---------------- ADMIN (front) ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="SSG DIGITAL ATTENDANCE")
SLIDE_MAP["digital_attendance_front"]=PAGE["n"]
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Certification / Credentials","20+ years of training and industry experience."),
  ("Delivers","WSQ courses on customer-centric branding, marketing communications and business skills."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with branding, marketing or communications (if any).",
 "A brand you personally admire — and why it works on you."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Actively participate — no question is too small.",
 "Respect each other's views: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required for funding eligibility."],
 kicker="HOUSEKEEPING",cols=2,size=15)
tile_grid("Skills Framework",[
 ("TSC Title",C.TSC_TITLE),
 ("TSC Code",C.TSC_CODE),
 ("Proficiency Level",C.TSC_LEVEL),
 ("Structure","4 Learning Units (LU1–LU4), each mapped to a Learning Outcome.")],
 kicker="COURSE ACCREDITATION",cols=2,size=15)
tile_grid("Course Outline",
 [(f"{t['code']} — {t['title']}", t["subtitle"]) for t in C.TOPICS],
 kicker="4 LEARNING UNITS",cols=2,size=14)
SLIDE_MAP["lesson_plan_slide"]=None
two_col(f"Lesson Plan — {C.DAYS} Days",
 [(f"Day 1 — {C.DAY_THEMES[1]}",0),
  ("LU1: Stakeholders and Organisation (Activities 1–4)",1),
  ("LU2: Customer Influence (Activities 5–8)",1),
  ("LU3: Branding in Marketing — begins (Activities 9–12)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
  ("LU3: Branding in Marketing — continued",1),
  ("LU4: Branding Effectiveness (Activities 13–17)",1),
  ("Assessment: WA (1h) + Case Study (1h), 4:00–6:00pm",1),
  ("Daily timing",0),("Day 1: 9:30am–6:30pm, 1-hour lunch. Day 2: 9:30am–6:00pm, 45-min lunch (to fit the assessment).",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2 & Assessment")
SLIDE_MAP["lesson_plan_slide"]=PAGE["n"]
tile_grid("Learning Outcomes",[
 ("Stakeholders & Organisation","Identify stakeholders/audiences and draft branding designs & reputation assessments."),
 ("Customer Influence","Capture and analyse customer insight through active listening."),
 ("Branding in Marketing","Execute branding campaigns, events and PR activities that build awareness."),
 ("Branding Effectiveness","Evaluate reputation and PR performance against KPIs and recommend improvements.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
SLIDE_MAP["briefing_assessment"]=PAGE["n"]
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
SLIDE_MAP["assessment_front"]=PAGE["n"]
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
lms_slide()

# ---------------- TOPICS + ACTIVITIES ----------------
# Each Learning Unit: section divider -> Key Concepts -> brand-trust insight
# slides -> per K/A sub-topic (divider, What is, How it Works, Compare,
# Supporting Data, optional Real-World Example) -> the matching in-class
# activity (scenario/role play -> discussion -> reflection) -> LU recap ->
# closing statement. Theory always precedes practice.
TOPIC_ACTS = {t["num"]: [a for a in ACT if a["topic"]==t["num"]] for t in C.TOPICS}
for t in C.TOPICS:
    section(f"{t['code']}", t["title"], t["code"], t["subtitle"])
    SLIDE_MAP[f"topic{t['num']}_section"]=PAGE["n"]
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"{t['code']} · KEY CONCEPTS", cols=2, size=14)
    for kind,spec in BRAND_TRUST.get(t["num"],[]):
        render_insight(kind,spec)
    # communication-tactics enrichment (charts, scorecards, process maps) drawn
    # from the industry references — the deck's "communication" half
    for kind,spec in COMMS.get(t["num"],[]):
        render_insight(kind,spec)
    acts=TOPIC_ACTS[t["num"]]
    for ti,a in enumerate(acts,1):
        t_tag=f"LO{t['num']} · {t['code']} · T{ti}"
        sub_divider(t_tag, a["t_statement"], f"T{ti}")
        render_what_is(a,t_tag)
        render_process(a,t_tag)
        render_compare(a,t_tag)
        render_visual(a,t_tag)
        if a.get("case_example_company"):
            case_study_callout(a["case_example_company"], a["case_example_headline"],
                                a["case_example_insight"], kicker=f"{t_tag} · REAL-WORLD EXAMPLE")
        act_kicker=f"{t_tag} · ACTIVITY {a['num']}"
        scenario_slide(act_kicker, a["title"], a["case_type"], a["case_scenario"], a["build"], a["duration"])
        SLIDE_MAP[f"activity{a['num']}"]=PAGE["n"]
        if a["case_type"]=="role_play" and a.get("roles"):
            role_cards_slide(act_kicker, a["title"], a["roles"])
        discussion_slide(act_kicker, a["title"], a["discussion_prompts"])
        reflection_slide(act_kicker, a["title"], a["reflection_points"], a["debrief_check"])
    recap_items=list({x["objective"]:x for x in acts}.values())[:6]
    tile_grid(f"Recap — {t['title']}",
              [(a["title"], f"You can now: {a['objective']}.") for a in recap_items],
              kicker=f"{t['code']} RECAP", cols=2, size=14)
    big_statement(t["closer_quote"], t["closer_attribution"], kicker=f"{t['code']} · WHY IT MATTERS", color=VIOLET)

playbook(C.COURSE_PLAYBOOK["title"],
         [(h,body) for _num,h,body in C.COURSE_PLAYBOOK["items"]],
         kicker="COURSE WRAP-UP", tagline=C.COURSE_PLAYBOOK["closing"])

# ---------------- CLOSE ----------------
# Digital Attendance must be the LAST admin slide before Thank You (WSQ hard rule) —
# nothing else goes after it, so the course-closer playbook slide above sits BEFORE
# this whole admin block, not after it.
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Stakeholders & Organisation","Mapped stakeholders and audiences; drafted brand designs and reputation assessments."),
 ("Customer Influence","Captured and documented customer perception through active listening."),
 ("Branding in Marketing","Executed branding and PR campaigns aligned to strategy, plan and budget."),
 ("Branding Effectiveness","Measured reputation and PR performance against KPIs and proposed improvements.")],
 kicker="LEARNING OUTCOMES",cols=2,size=15)
tile_grid("Assessment",[
 ("Written Assessment (WA)","SAQ — 1 hour, open book."),
 ("Case Study (CS)","1 hour, open book."),
 ("Open Book","Slides, Learner Guide and approved materials only."),
 ("Submit","On the LMS — lms-tms.tertiaryinfotech.com")],
 kicker="WRAP-UP",cols=2,size=14)
SLIDE_MAP["assessment_end"]=PAGE["n"]
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="SSG DIGITAL ATTENDANCE")
SLIDE_MAP["digital_attendance_end"]=PAGE["n"]
big_statement("Thank You!","You are now equipped to build customer-centric brand communication that earns trust and drives results.",
              "SEE YOU AT THE NEXT ONE",color=TEAL)

# ---------------- motion: transitions + entrance animations ----------------
# Transitions are chosen by slide ROLE rather than applied uniformly, so the
# movement carries meaning: a push means "new section", a wipe means "new task",
# a zoom marks a deliberate beat, and ordinary content slides fade quietly.
# Roles are inferred from each slide's own content, which keeps this independent
# of the order the builders above happen to run in.

def _classify(sl, idx, total):
    """Infer a slide's role for transition selection."""
    texts=[sh.text_frame.text for sh in sl.shapes if sh.has_text_frame]
    joined=" ".join(texts)
    kicker=texts[0] if texts else ""
    if idx==0: return "cover"
    if idx>=total-1: return "closing"
    # section dividers carry a big ghosted LU/T number and a short body
    if any(t.strip() in (f"LU{n}" for n in range(1,9)) for t in texts): return "section"
    if any(re.fullmatch(r"T\d+", t.strip()) for t in texts): return "subdivider"
    if "· ACTIVITY" in joined or "THE SCENARIO" in joined.upper(): return "scenario"
    if "YOUR ROLES" in joined.upper() or "ROLE CARDS" in joined.upper(): return "roles"
    if "DISCUSSION" in joined.upper() or "DECISION PROMPTS" in joined.upper(): return "discussion"
    if "REFLECT" in joined.upper() or "DEBRIEF" in joined.upper(): return "reflection"
    if "WHY IT MATTERS" in joined.upper() or "SEE YOU AT" in joined.upper(): return "big"
    if any(k in joined for k in ("Digital Attendance","Assessment Flow","Briefing for Assessment",
                                 "Ground Rules","Lesson Plan","About the Trainer")): return "admin"
    return "content"

import re
_total=len(prs.slides._sldIdLst)
_roles=[_classify(sl,i,_total) for i,sl in enumerate(prs.slides)]
_plan=motion.build_plan(_roles)
_tally=motion.apply_transitions(prs,_plan,speed="fast")

# Entrance animations: on card/tile slides, reveal the cards in reading order so
# the trainer can talk to each one. Only the body shapes animate — the header,
# rule and footer are present from the start.
def _animate_cards(sl, role):
    if role in ("cover","closing","big","section"): return
    body=[sh for sh in sl.shapes
          if sh.has_text_frame and sh.top is not None and sh.top > Inches(2.1)
          and sh.height is not None and sh.height > Inches(0.5)]
    body.sort(key=lambda s:(round(s.top/Inches(0.5)), s.left))
    if 2 <= len(body) <= 8:
        motion.animate(sl, body, dur=350, stagger=True)

for sl,role in zip(prs.slides,_roles):
    _animate_cards(sl,role)

print("Transitions:", ", ".join(f"{k}×{v}" for k,v in sorted(_tally.items())))

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
with open(os.path.join(HERE,"slide_map.json"),"w") as f:
    json.dump(SLIDE_MAP,f,indent=2)
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
print("Saved slide_map.json")
